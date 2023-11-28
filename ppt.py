from pptx import Presentation 
from pptx.util import Inches
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from openpyxl import load_workbook

def getPPT():
    file_name = "E:\\001_QA_Generator\\input.xlsx"

    load_wb = load_workbook(file_name, data_only=True)
    load_ws = load_wb['Sheet1']
    row_start=1
    row_count = load_ws.max_row
    i=1

    slide_number = 1
    X = Presentation()
    X.slide_width = Inches(16)
    X.slide_height = Inches(9)
    Second_Layout = X.slide_layouts[6]

    while(i<=row_count):
        part_number=load_ws.cell(row=row_start,column=1).value
        heading=load_ws.cell(row=row_start,column=2).value
        second_slide = X.slides.add_slide(Second_Layout)
        logo=second_slide.shapes.add_picture("E:\\001_QA_Generator\\logo.png",Inches(13.8), Inches(0.05),width=Inches(2.07),height = Inches(0.77))
        #------------------------heading-------------------------#
        textbox = second_slide.shapes.add_textbox(Inches(1.5), Inches(0.15),Inches(11), Inches(0.5))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = str(heading)
        paragraph.font.size=Pt(38)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 109, 117)    
        #------------------slide no-------------------#
        textbox = second_slide.shapes.add_textbox(Inches(0.5), Inches(1.5),Inches(0.5), Inches(1.5))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = str(slide_number)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255,255,255)    
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(0, 109, 117)
        #-------------------JSON----------------#
        png_name=part_number+"_json.png"
        image=second_slide.shapes.add_picture(f"E:\\001_QA_Generator\\png_images\\{part_number}\\{png_name}", Inches(1.25), Inches(2.2),height = Inches(2),width=Inches(2))
        textbox = second_slide.shapes.add_textbox(Inches(1.75), Inches(1.25),Inches(0.2), Inches(0.2))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = "JSON"
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 109, 117)

        #-------------SIMP_STP-------------------#
        png_name=part_number+"_simp.png"
        image=second_slide.shapes.add_picture(f"E:\\001_QA_Generator\\png_images\\{part_number}\\{png_name}", Inches(4.75), Inches(2.2),height = Inches(2),width=Inches(2))
        textbox = second_slide.shapes.add_textbox(Inches(5), Inches(1.25),Inches(0.2), Inches(0.2))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = "SIMP_STP"
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 109, 117)
        #-----------------SWX--------------------#
        png_name=part_number+"_sldprt.png"
        image=second_slide.shapes.add_picture(f"E:\\001_QA_Generator\\png_images\\{part_number}\\{png_name}", Inches(1.25), Inches(5.5),height = Inches(2),width=Inches(2))
        textbox = second_slide.shapes.add_textbox(Inches(1.75), Inches(4.5),Inches(0.2), Inches(0.2))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = "SWX"
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 109, 117)
        #------------------PRT-------------------#
        png_name=part_number+"_prt.png"
        image=second_slide.shapes.add_picture(f"E:\\001_QA_Generator\\png_images\\{part_number}\\{png_name}", Inches(4.75), Inches(5.5),height = Inches(2),width=Inches(2))
        textbox = second_slide.shapes.add_textbox(Inches(5.25), Inches(4.5),Inches(0.2), Inches(0.2))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = "PRT"
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 109, 117)
        #---------------2dpdf--------------------#
        pdf_name=part_number+"_pdf.png"
        image=second_slide.shapes.add_picture(f"E:\\001_QA_Generator\\png_images\\{part_number}\\{pdf_name}", Inches(7.5), Inches(2.1),height = Inches(5.5),width=Inches(8))
        textbox = second_slide.shapes.add_textbox(Inches(11), Inches(1.25),Inches(0.2), Inches(0.2))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = "2dpdf"
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 109, 117)
        #-------------------bottom link------------
        textbox = second_slide.shapes.add_textbox(Inches(1.25), Inches(8),Inches(0.3), Inches(11))
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = f"https://qa.product-config.net/catalog3/d/grainger/?c=products&cid=root&id={part_number}"
        paragraph.font.size=Pt(18)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 109, 117)

        slide_number+=1
        row_start+=1
        i+=1
    X.save("E:\\001_QA_Generator\\First_presentation.pptx")

if __name__ == "__main__":
    getPPT()
